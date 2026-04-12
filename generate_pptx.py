import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, top_inches=0.3):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top_inches), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.name = 'Segoe UI'
    p.font.color.rgb = RGBColor(0, 200, 255) # Cyan highlight
    return title_box

def add_content(slide, text, left, top, width, height, font_size=16, bold=False, color=RGBColor(255, 255, 255)):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = 'Segoe UI'
    p.font.bold = bold
    p.font.color.rgb = color
    return box, p

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625) # 16:9 ratio

    BG_COLOR = RGBColor(18, 18, 18)
    TEXT_COLOR = RGBColor(240, 240, 240)
    HIGHLIGHT_1 = RGBColor(0, 255, 170) # Neon green
    HIGHLIGHT_2 = RGBColor(255, 85, 85) # Neon red
    HIGHLIGHT_3 = RGBColor(255, 204, 0) # Neon yellow

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    set_slide_background(slide1, BG_COLOR)
    
    add_content(slide1, "🤖 AI CODING TOOLS\nCLAUDE CODE vs GEMINI CLI vs ANTIGRAVITY", 
                1, 1.5, 8, 2, font_size=36, bold=True, color=HIGHLIGHT_1)
    add_content(slide1, "Hướng dẫn dành cho người mới bắt đầu • Data Analyst Team • 2026", 
                1, 3.2, 8, 1, font_size=18, color=TEXT_COLOR)

    # --- Slide 2: 3 Công cụ AI ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2, BG_COLOR)
    add_title(slide2, "PHẦN 1 — 3 CÔNG CỤ AI LÀ GÌ?")
    add_content(slide2, "Trong thế giới phát triển phần mềm, ba trợ lý AI CLI đang nổi bật nhất hiện nay hỗ trợ DA làm việc hiệu quả theo những cách khác nhau.",
                0.5, 1.2, 9, 1, font_size=16, color=TEXT_COLOR)
    
    # boxes
    y_off = 2.0
    add_content(slide2, "🟢 Claude Code (Anthropic)", 0.5, y_off, 2.8, 0.5, 18, True, HIGHLIGHT_1)
    add_content(slide2, "Agent tự hành. Lập luận độc lập cực tốt, ít cần can thiệp. Phụ thuộc vào file CLAUDE.md để ghi nhớ nội quy.", 0.5, y_off+0.5, 2.8, 2, 14, False, TEXT_COLOR)

    add_content(slide2, "🔵 Gemini CLI (Google)", 3.6, y_off, 2.8, 0.5, 18, True, RGBColor(50, 150, 255))
    add_content(slide2, "Gọn nhẹ, scripting & API. Đọc file GEMINI.md tại chỗ, giải quyết tức thì, không có bộ nhớ dài hạn chuyên biệt.", 3.6, y_off+0.5, 2.8, 2, 14, False, TEXT_COLOR)

    add_content(slide2, "🟣 AntiGravity (Google DeepMind)", 6.7, y_off, 2.8, 0.5, 18, True, RGBColor(200, 100, 255))
    add_content(slide2, "Hybrid Interface + Autonomous Agent. Tự học và nhớ qua Knowledge Items. Có Browser Subagent và hệ thống báo cáo Artifacts mượt mà.", 6.7, y_off+0.5, 2.8, 2, 14, False, TEXT_COLOR)

    # --- Slide 3: So Sánh ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3, BG_COLOR)
    add_title(slide3, "PHẦN 2 — SO SÁNH TÍNH NĂNG CHÍNH")
    
    shapes = slide3.shapes
    rows, cols = 5, 4
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(0.4)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # style table cols
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.3)
    table.columns[2].width = Inches(2.3)
    table.columns[3].width = Inches(2.4)

    data = [
        ["Tiêu Chí", "Claude Code", "Gemini CLI", "AntiGravity"],
        ["Bộ nhớ dài hạn\n(Long-term memory)", "🔶 Thủ công\nPhải ghi file", "❌ Không có", "✅ Tự động qua KIs"],
        ["Trình bày tự động", "❌ Text trên terminal", "❌ Text trên terminal", "✅ Markdown, Slide, \nBiểu đồ trực quan"],
        ["Điều khiển trình duyệt", "❌ Không tích hợp", "❌ Không có", "✅ Browser Subagent"],
        ["Hỗ trợ MCP", "✅ Rất mạnh (Anthropic chuẩn)", "🔶 Hạn chế", "✅ Có hỗ trợ"]
    ]

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = 'Segoe UI'
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = HIGHLIGHT_1
                else:
                    p.font.color.rgb = TEXT_COLOR
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

    # --- Slide 4: AntiGravity Features ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4, BG_COLOR)
    add_title(slide4, "PHẦN 3 — 4 TÍNH NĂNG ĐỈNH CAO CỦA ANTIGRAVITY")
    
    add_content(slide4, "1. Autonomous Agent", 0.5, 1.2, 4.0, 0.5, 16, True, HIGHLIGHT_3)
    add_content(slide4, "Tự lập kế hoạch, chạy lệnh, đọc/ghi file, tạo ảnh báo cáo mà không cần can thiệp.", 0.5, 1.6, 4.0, 1, 14, False, TEXT_COLOR)
    
    add_content(slide4, "2. Bộ não Knowledge Items (KIs)", 0.5, 3.2, 4.0, 0.5, 16, True, HIGHLIGHT_3)
    add_content(slide4, "Học từ thói quen cũ. Lưu naming convention, cách sửa lỗi vào KIs để dùng tự động cho session sau.", 0.5, 3.6, 4.0, 1, 14, False, TEXT_COLOR)

    add_content(slide4, "3. Browser Subagent", 5.0, 1.2, 4.5, 0.5, 16, True, HIGHLIGHT_3)
    add_content(slide4, "Điều khiển trình duyệt ngầm, click, cuộn, thu thập Data, ghi video webp hoàn toàn tự động.", 5.0, 1.6, 4.5, 1, 14, False, TEXT_COLOR)

    add_content(slide4, "4. Báo cáo Artifacts Chuyên nghiệp", 5.0, 3.2, 4.5, 0.5, 16, True, HIGHLIGHT_3)
    add_content(slide4, "Render báo cáo Markdown đẹp mắt, bảng biểu, sơ đồ Mermaid, Carousel. Giúp handover tuyệt vời.", 5.0, 3.6, 4.5, 1, 14, False, TEXT_COLOR)


    # --- Slide 5: MCP & Tips ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5, BG_COLOR)
    add_title(slide5, "PHẦN 4 — MCP LÀ GÌ? TIPS THỰC CHIẾN")
    
    add_content(slide5, "🔌 MCP (Model Context Protocol)", 0.5, 1.2, 9.0, 0.5, 18, True, RGBColor(255,0,255))
    add_content(slide5, "Chuẩn kết nối mở (như ổ cắm đa năng). Giúp AI cắm vào BigQuery, GitHub, Notion, Slack.\nVí dụ: 'Truy vấn BigQuery doanh thu T3 và gửi report qua Notion/Slack' → AI tự làm từ A-Z.", 0.5, 1.7, 9.0, 1, 14, False, TEXT_COLOR)

    add_content(slide5, "💡 TIPS CHO DATALYS TEAM", 0.5, 3.0, 9.0, 0.5, 18, True, HIGHLIGHT_3)
    add_content(slide5, "• Prompt thẩm mỹ: Thêm 'Premium design, Modern typography' để output đẹp.\n• Kích hoạt KIs: Chủ động nhắc AI lưu pattern vào KIs.\n• _GEMINI.md cục bộ: Hướng dẫn riêng cho project.\n• Built-in model: Chạy tốt nhất với model tích hợp của hệ thống Agent.", 0.5, 3.5, 9.0, 2, 14, False, TEXT_COLOR)

    # --- Slide 6: NEW - CHATBOT INSIGHT --- 
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide6, BG_COLOR)
    add_title(slide6, "🎯 ỨNG DỤNG MỚI: CHATBOT INSIGHT TRÊN GOOGLE CHAT", top_inches=0.2)
    
    # Content block
    add_content(slide6, "Hệ Thống:", 0.5, 1.2, 2.0, 0.4, 16, True, HIGHLIGHT_1)
    add_content(slide6, "Sử dụng Google Pub/Sub (host local 24/7 chờ kết nối).", 2.0, 1.2, 7.5, 0.4, 16, False, TEXT_COLOR)
    
    add_content(slide6, "Concept Hoạt Động (User Workflow):", 0.5, 1.8, 9.0, 0.4, 16, True, HIGHLIGHT_2)
    
    # Steps
    add_content(slide6, "1️⃣ Query (Truy vấn):", 0.5, 2.4, 2.5, 0.4, 15, True, rgb_color(0, 200, 255))
    add_content(slide6, "Người dùng tag tên Bot trên group chat.\nVD: @BOT_INSIGHT: Biến động doanh thu tháng này", 2.8, 2.4, 6.7, 0.6, 14, False, TEXT_COLOR)

    add_content(slide6, "2️⃣ Acknowledge:", 0.5, 3.1, 2.5, 0.4, 15, True, rgb_color(0, 200, 255))
    add_content(slide6, "Bot lắng nghe, lập tức phản hồi 'Đợi tí nhé sếp...' để ghi nhận request.", 2.8, 3.1, 6.7, 0.4, 14, False, TEXT_COLOR)

    add_content(slide6, "3️⃣ Process:", 0.5, 3.7, 2.5, 0.4, 15, True, rgb_color(0, 200, 255))
    add_content(slide6, "Bot tự động query data từ Dremio, phân tích và trích xuất Insight.", 2.8, 3.7, 6.7, 0.4, 14, False, TEXT_COLOR)

    add_content(slide6, "4️⃣ Feedback:", 0.5, 4.3, 2.5, 0.4, 15, True, rgb_color(0, 200, 255))
    add_content(slide6, "Bot tổng hợp báo cáo và trả lời trực tiếp trong chính thread đó trên Google Chat.", 2.8, 4.3, 6.7, 0.5, 14, False, TEXT_COLOR)

    prs.save("ANTIGRAVITY_FIX_v2.pptx")

def rgb_color(r, g, b):
    return RGBColor(r, g, b)

if __name__ == '__main__':
    main()
