import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

prs = Presentation('ANTIGRAVITY_FIX_v2.pptx')

for i, slide in enumerate(prs.slides):
    print(f"=== SLIDE {i+1} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if texts:
                print(f"  [{shape.name}]")
                for t in texts:
                    print(f"    {t}")
    print()
